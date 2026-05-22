"""PPTX text extraction using python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from parsers.base import BaseParser, ParserError


class PPTXParser(BaseParser):
    extensions = (".pptx", ".ppt")

    def parse(self, file_path: Path) -> str:
        try:
            presentation = Presentation(str(file_path))
            parts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text.strip())
            return "\n".join(parts)
        except PackageNotFoundError as exc:
            raise ParserError("Invalid or corrupted PPTX file.") from exc
        except Exception as exc:
            raise ParserError(f"Failed to parse PPTX: {exc}") from exc
